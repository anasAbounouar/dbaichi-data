#!/usr/bin/env python3
"""
Add all images from HTML presentation to the PPTX file - Version 2
This version properly embeds images even if slides already have some content
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt

def parse_html_for_images():
    """Parse HTML to extract slide-to-image mappings with better detection"""
    with open('/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation.html', 'r') as f:
        html_content = f.read()

    # Split by slide divs more accurately
    slide_sections = html_content.split('<div class="slide')

    image_mappings = []

    for slide_num, section in enumerate(slide_sections[1:], 1):  # Skip first split (before slides)
        # Find the end of this slide
        slide_end = section.find('<div class="slide')
        if slide_end == -1:
            slide_end = section.find('</div>')

        slide_content = section[:slide_end] if slide_end != -1 else section

        # Find images in this slide
        img_pattern = r'<img[^>]+src="([^"]+)"[^>]*alt="([^"]*)"'
        images = re.findall(img_pattern, slide_content)

        if images:
            # Determine if full-image slide
            is_full_image = 'full-image-slide' in slide_content

            for img_src, img_alt in images:
                # Clean up the path
                img_path = img_src.replace('images/', '')

                # Skip logos on title slide (they're decorative)
                if 'logo' in img_path.lower() and slide_num <= 2:
                    continue

                image_mappings.append({
                    'slide_num': slide_num,
                    'image_path': img_path,
                    'is_full_image': is_full_image,
                    'alt_text': img_alt
                })

    return image_mappings

def add_images_to_pptx(pptx_path, image_mappings, images_dir):
    """Add images to PPTX slides - force add all images"""
    prs = Presentation(pptx_path)

    print(f"\n{'='*60}")
    print(f"Adding images to PPTX")
    print(f"{'='*60}\n")

    images_added = 0
    images_skipped = 0

    for mapping in image_mappings:
        slide_num = mapping['slide_num']
        image_path = mapping['image_path']
        is_full_image = mapping['is_full_image']
        alt_text = mapping.get('alt_text', '')

        # Adjust slide number (0-indexed in python-pptx)
        slide_idx = slide_num - 1

        if slide_idx >= len(prs.slides):
            print(f"⚠ Slide {slide_num} not found in PPTX, skipping")
            images_skipped += 1
            continue

        slide = prs.slides[slide_idx]

        # Full path to image
        full_image_path = os.path.join(images_dir, image_path)

        if not os.path.exists(full_image_path):
            print(f"⚠ Image not found: {image_path}, skipping")
            images_skipped += 1
            continue

        try:
            # Determine image placement based on slide type
            if is_full_image:
                # Full slide image - large and centered
                # Remove existing images first to avoid clutter
                shapes_to_remove = []
                for shape in slide.shapes:
                    if shape.shape_type == 13:  # Picture
                        shapes_to_remove.append(shape)

                for shape in shapes_to_remove:
                    sp = shape.element
                    sp.getparent().remove(sp)

                # Add new full-size image
                left = Inches(0.5)
                top = Inches(1.0)
                height = Inches(6.5)
                slide.shapes.add_picture(full_image_path, left, top, height=height)
                print(f"✓ Added FULL image to slide {slide_num}: {image_path}")
                images_added += 1
            else:
                # Check if this specific image already exists (check by looking at text content)
                # For now, just add it

                # Position based on slide layout
                # If slide has bullet points, place on right side
                # If slide is empty, center it

                has_text_content = any(
                    shape.has_text_frame and shape.text.strip()
                    for shape in slide.shapes
                )

                if has_text_content:
                    # Place on right side
                    left = Inches(7.5)
                    top = Inches(2.0)
                    height = Inches(4.0)
                else:
                    # Center it
                    left = Inches(2.5)
                    top = Inches(2.0)
                    height = Inches(4.5)

                pic = slide.shapes.add_picture(full_image_path, left, top, height=height)
                print(f"✓ Added image to slide {slide_num}: {image_path}")
                images_added += 1

        except Exception as e:
            print(f"✗ Error adding image {image_path} to slide {slide_num}: {e}")
            images_skipped += 1

    # Save the modified PPTX
    output_path = pptx_path.replace('_converted.pptx', '_with_all_images.pptx')
    prs.save(output_path)

    print(f"\n{'='*60}")
    print(f"Summary")
    print(f"{'='*60}")
    print(f"✅ Images added: {images_added}")
    print(f"⚠️  Images skipped: {images_skipped}")
    print(f"📁 Output file: {output_path}")
    print(f"📊 File size: {os.path.getsize(output_path) / (1024*1024):.1f} MB")
    print(f"{'='*60}\n")

    return output_path

def main():
    base_dir = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle'
    pptx_path = os.path.join(base_dir, 'presentation_converted.pptx')
    images_dir = os.path.join(base_dir, 'images')

    print("📖 Step 1: Parsing HTML to find images...")
    image_mappings = parse_html_for_images()
    print(f"   Found {len(image_mappings)} image references in HTML\n")

    print("🖼️  Step 2: Embedding images into PPTX...")
    output_path = add_images_to_pptx(pptx_path, image_mappings, images_dir)

    print("\n✅ COMPLETE! Your PPTX now has all images embedded.")
    print(f"📁 Open: {output_path}")
    print("\n💡 The images are now embedded so you can share this file anywhere!")

if __name__ == "__main__":
    main()
