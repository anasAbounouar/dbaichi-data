#!/usr/bin/env python3
"""
High-quality HTML to PowerPoint converter
Creates professional presentation with proper formatting
"""
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from lxml import html as lxml_html
from pptx.enum.shapes import MSO_SHAPE

# Read and parse HTML
html_file = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation.html'
with open(html_file, 'r', encoding='utf-8') as f:
    content = f.read()

tree = lxml_html.fromstring(content)

# Initialize presentation with widescreen format
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

# Define colors from the original design
TEAL_COLOR = RGBColor(20, 83, 95)  # #14535F
BLUE_COLOR = RGBColor(59, 130, 246)
GREEN_COLOR = RGBColor(34, 197, 94)
ORANGE_COLOR = RGBColor(251, 146, 60)
RED_COLOR = RGBColor(239, 68, 68)
YELLOW_COLOR = RGBColor(250, 204, 21)
GRAY_DARK = RGBColor(31, 41, 55)
GRAY_MED = RGBColor(107, 114, 128)
GRAY_LIGHT = RGBColor(243, 244, 246)

def clean_text(text):
    """Clean and normalize text"""
    if text is None:
        return ""
    text = re.sub(r'\s+', ' ', text.strip())
    # Remove emoji artifacts
    text = re.sub(r'[●□■▪▫]', '', text)
    return text

def add_background(slide, color=None, is_gradient=False):
    """Add colored background to slide"""
    background = slide.background
    fill = background.fill
    if is_gradient:
        fill.gradient()
        fill.gradient_stops[0].color.rgb = TEAL_COLOR
        fill.gradient_stops[1].color.rgb = RGBColor(26, 104, 118)
    elif color:
        fill.solid()
        fill.fore_color.rgb = color
    else:
        fill.solid()
        fill.fore_color.rgb = RGBColor(255, 255, 255)

def create_title_slide(title, subtitle, author, supervisors, year):
    """Create professional title slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
    add_background(slide, is_gradient=True)

    # Add logos (if images exist)
    logo_oracle = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/images/logo_oracle.png'
    logo_ehtp = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/images/logo_ehtp.jpg'

    if os.path.exists(logo_oracle):
        slide.shapes.add_picture(logo_oracle, Inches(0.5), Inches(0.5), height=Inches(0.8))
    if os.path.exists(logo_ehtp):
        slide.shapes.add_picture(logo_ehtp, Inches(12), Inches(0.5), height=Inches(1))

    # Main title
    title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(11.333), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = title
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(48)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Subtitle
    if subtitle:
        subtitle_box = slide.shapes.add_textbox(Inches(1.5), Inches(4), Inches(10.333), Inches(0.8))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_frame.word_wrap = True
        for paragraph in subtitle_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(20)
            paragraph.font.color.rgb = RGBColor(191, 219, 254)

    # Author info
    if author:
        author_box = slide.shapes.add_textbox(Inches(1), Inches(5.5), Inches(3.5), Inches(1))
        author_frame = author_box.text_frame
        author_frame.text = f"Presented By:\n{author}"
        for paragraph in author_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Supervisors
    if supervisors:
        sup_box = slide.shapes.add_textbox(Inches(9), Inches(5.5), Inches(3.5), Inches(1))
        sup_frame = sup_box.text_frame
        sup_frame.text = f"Supervisors:\n{supervisors}"
        for paragraph in sup_frame.paragraphs:
            paragraph.font.size = Pt(14)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Year
    if year:
        year_box = slide.shapes.add_textbox(Inches(4), Inches(6.5), Inches(5.333), Inches(0.5))
        year_frame = year_box.text_frame
        year_frame.text = year
        for paragraph in year_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = RGBColor(255, 255, 255)

def create_section_divider(section_title):
    """Create section divider slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, is_gradient=True)

    # Section title
    title_box = slide.shapes.add_textbox(Inches(2), Inches(2.5), Inches(9.333), Inches(2))
    title_frame = title_box.text_frame
    title_frame.text = section_title
    title_frame.word_wrap = True
    for paragraph in title_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(60)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(5.5), Inches(4.8), Inches(2.333), Inches(0.08)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = YELLOW_COLOR
    line.line.fill.background()

def create_content_slide(title, content_data):
    """Create content slide with proper formatting"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_background(slide, RGBColor(255, 255, 255))

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.333), Inches(0.8))
    title_frame = title_box.text_frame
    title_frame.text = title
    for paragraph in title_frame.paragraphs:
        paragraph.font.size = Pt(36)
        paragraph.font.bold = True
        paragraph.font.color.rgb = GRAY_DARK

    return slide

def add_bullet_points(slide, items, left, top, width, height, font_size=16):
    """Add formatted bullet points to slide"""
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for idx, item in enumerate(items):
        if idx > 0:
            text_frame.add_paragraph()
        p = text_frame.paragraphs[idx]
        p.text = item
        p.level = 0
        p.font.size = Pt(font_size)
        p.font.color.rgb = GRAY_DARK
        p.space_after = Pt(8)

def add_stat_boxes(slide, stats, top_position):
    """Add colored statistic boxes"""
    num_stats = len(stats)
    if num_stats == 0:
        return

    box_width = Inches(2.2)
    spacing = Inches(0.2)
    total_width = (box_width * num_stats) + (spacing * (num_stats - 1))
    start_left = (Inches(13.333) - total_width) / 2

    colors = [
        RGBColor(219, 234, 254),  # blue-100
        RGBColor(220, 252, 231),  # green-100
        RGBColor(243, 232, 255),  # purple-100
        RGBColor(254, 235, 200),  # orange-100
        RGBColor(254, 226, 226),  # red-100
    ]

    for idx, stat in enumerate(stats[:5]):
        left_pos = start_left + (idx * (box_width + spacing))

        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left_pos, top_position, box_width, Inches(1.2)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = colors[idx % len(colors)]
        shape.line.fill.background()

        # Add text
        text_frame = shape.text_frame
        text_frame.clear()

        # Value
        p1 = text_frame.paragraphs[0]
        p1.text = stat.get('value', '')
        p1.alignment = PP_ALIGN.CENTER
        p1.font.size = Pt(24)
        p1.font.bold = True
        p1.font.color.rgb = TEAL_COLOR

        # Label
        p2 = text_frame.add_paragraph()
        p2.text = stat.get('label', '')
        p2.alignment = PP_ALIGN.CENTER
        p2.font.size = Pt(14)
        p2.font.color.rgb = GRAY_MED

# Extract slides from HTML
slides_elements = tree.xpath('//div[@class="slide"]')
print(f"Found {len(slides_elements)} slides in HTML")

# Process each slide
for idx, slide_elem in enumerate(slides_elements):
    print(f"Processing slide {idx + 1}...")

    try:
        # Check if it's a divider slide
        is_divider = slide_elem.get('data-divider') == 'true'

        # Get title
        h2_elements = slide_elem.xpath('.//h2')
        h1_elements = slide_elem.xpath('.//h1')

        if idx == 0:
            # Title slide
            title_elem = h1_elements[0] if h1_elements else None
            subtitle_elem = slide_elem.xpath('.//p[@class="text-2xl"]')
            author_elem = slide_elem.xpath('.//div[contains(., "Presented By")]/following-sibling::p')
            supervisor_elem = slide_elem.xpath('.//div[contains(., "Supervisors")]/following-sibling::p')
            year_elem = slide_elem.xpath('.//p[contains(., "Academic year")]')

            title = clean_text(title_elem.text_content()) if title_elem is not None else "Presentation"
            subtitle = clean_text(subtitle_elem[0].text_content()) if subtitle_elem else ""
            author = clean_text(author_elem[0].text_content()) if author_elem else ""
            supervisors = "\n".join([clean_text(p.text_content()) for p in supervisor_elem[:2]])
            year = clean_text(year_elem[0].text_content()) if year_elem else ""

            create_title_slide(title, subtitle, author, supervisors, year)

        elif is_divider:
            # Section divider
            title = clean_text(h2_elements[0].text_content()) if h2_elements else f"Section {idx}"
            create_section_divider(title)

        else:
            # Regular content slide
            title = clean_text(h2_elements[0].text_content()) if h2_elements else f"Slide {idx + 1}"
            slide = create_content_slide(title, {})

            # Extract images
            img_elements = slide_elem.xpath('.//img[contains(@src, "images/")]')
            has_large_image = False

            for img_elem in img_elements:
                img_src = img_elem.get('src')
                if img_src and 'logo' not in img_src.lower():
                    img_path = os.path.join('/Users/anasabounouar/Downloads/dbaichi/pfe-oracle', img_src)
                    if os.path.exists(img_path):
                        try:
                            # Check if image should be constrained
                            if 'factory_pattern' in img_src or 'dependency_resolution' in img_src:
                                slide.shapes.add_picture(img_path, Inches(3.5), Inches(1.5), width=Inches(6.5))
                            else:
                                slide.shapes.add_picture(img_path, Inches(1), Inches(1.5), width=Inches(11.333))
                            has_large_image = True
                        except Exception as e:
                            print(f"  Warning: Could not add image {img_src}: {e}")

            # Extract stat boxes (colored boxes with numbers)
            stat_divs = slide_elem.xpath('.//div[contains(@class, "bg-blue-50") or contains(@class, "bg-green-50") or contains(@class, "bg-purple-50")]')
            if len(stat_divs) >= 3 and not has_large_image:
                stats = []
                for stat_div in stat_divs[:5]:
                    value_elem = stat_div.xpath('.//div[contains(@class, "font-bold")]')
                    label_elem = stat_div.xpath('.//p[contains(@class, "font-medium")]')
                    if value_elem and label_elem:
                        stats.append({
                            'value': clean_text(value_elem[0].text_content()),
                            'label': clean_text(label_elem[0].text_content())
                        })
                if stats:
                    add_stat_boxes(slide, stats, Inches(2))

            # Extract bullet points
            li_elements = slide_elem.xpath('.//li')
            if li_elements and not has_large_image:
                bullets = []
                for li in li_elements[:10]:
                    text = clean_text(li.text_content())
                    if text and len(text) > 3:
                        bullets.append(text)

                if bullets:
                    top_pos = Inches(5.5) if stat_divs else Inches(1.8)
                    add_bullet_points(slide, bullets, Inches(0.8), top_pos, Inches(11.5), Inches(4.5), 14)

            # Extract grid content boxes
            if not has_large_image and not li_elements and not stat_divs:
                content_boxes = slide_elem.xpath('.//div[contains(@class, "grid")]//div[contains(@class, "bg-")]')
                if content_boxes:
                    content_items = []
                    for box in content_boxes[:6]:
                        h3 = box.xpath('.//h3 | .//h4')
                        p = box.xpath('.//p')
                        if h3:
                            title_text = clean_text(h3[0].text_content())
                            desc_text = " ".join([clean_text(par.text_content()) for par in p[:2]])
                            if title_text:
                                content_items.append(f"{title_text}: {desc_text}" if desc_text else title_text)

                    if content_items:
                        add_bullet_points(slide, content_items, Inches(0.8), Inches(1.8), Inches(11.5), Inches(5), 14)

    except Exception as e:
        print(f"  Error processing slide {idx + 1}: {e}")
        # Create a simple slide with just the title
        try:
            h2_elements = slide_elem.xpath('.//h2')
            title = clean_text(h2_elements[0].text_content()) if h2_elements else f"Slide {idx + 1}"
            create_content_slide(title, {})
        except:
            pass

# Save presentation
output_file = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation_quality.pptx'
prs.save(output_file)
print(f"\n✓ High-quality presentation saved to: {output_file}")
print(f"✓ Total slides created: {len(prs.slides)}")
