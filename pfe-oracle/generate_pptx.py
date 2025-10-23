#!/usr/bin/env python3
"""
Perfect HTML to PowerPoint Converter
Generates pixel-perfect PPTX matching the HTML presentation exactly
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Exact color palette from HTML
TEAL_PRIMARY = RGBColor(20, 83, 95)      # #14535F
TEAL_DARK = RGBColor(10, 57, 64)         # #0a3940
TEAL_LIGHT = RGBColor(15, 118, 110)      # #0f766e
RED_ACCENT = RGBColor(199, 70, 52)       # #C74634
LIGHT_GRAY_BG = RGBColor(248, 250, 252)  # #f8fafc
WHITE = RGBColor(255, 255, 255)
TEXT_GRAY = RGBColor(51, 65, 85)         # #334155
TEXT_LIGHT = RGBColor(100, 116, 139)     # #64748b

# Exact dimensions for 16:9 widescreen
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(5.625)

class PresentationBuilder:
    def __init__(self, images_dir):
        self.prs = Presentation()
        self.prs.slide_width = SLIDE_WIDTH
        self.prs.slide_height = SLIDE_HEIGHT
        self.images_dir = images_dir

    def add_title_slide(self):
        """Slide 1: Title slide with logos"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Dark teal gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = TEAL_PRIMARY
        fill.gradient_stops[1].color.rgb = TEAL_DARK

        # Logos at top
        logo_top = Inches(0.8)
        logo1_path = os.path.join(self.images_dir, "logo_ehtp.jpg")
        logo2_path = os.path.join(self.images_dir, "logo_graalvm.png")

        if os.path.exists(logo1_path):
            slide.shapes.add_picture(logo1_path, Inches(3.0), logo_top, height=Inches(0.8))
        if os.path.exists(logo2_path):
            slide.shapes.add_picture(logo2_path, Inches(6.0), logo_top, height=Inches(0.8))

        # Main title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(9), Inches(1.2))
        title_frame = title_box.text_frame
        title_frame.text = "Infrastructure as Code Support\nin Graal CI"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(44)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.3), Inches(9), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "JSON-Driven IaC Framework with Pulumi & GitLab CI/CD"
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(20)
        subtitle_frame.paragraphs[0].font.color.rgb = WHITE

        # Author info
        author_box = slide.shapes.add_textbox(Inches(0.5), Inches(4.2), Inches(9), Inches(1))
        author_frame = author_box.text_frame
        author_text = "Yassine DBAICHI\nÉcole Hassania des Travaux Publics\nOracle Labs - GraalVM RISQ Team\nAcademic Year 2024-2025"
        author_frame.text = author_text
        for paragraph in author_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(16)
            paragraph.font.color.rgb = WHITE
            paragraph.line_spacing = 1.8

    def add_agenda_slide(self):
        """Slide 2: Table of Contents with 2x3 card grid"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        # Title with red underline
        self._add_content_title(slide, "Agenda")

        # 2x3 card grid
        cards_data = [
            ("1. Company & Team Context", "Oracle, Oracle Labs, GraalVM RISQ Team mission and focus"),
            ("2. Problem & Solution Overview", "Current challenges, requirements, and proposed JSON-driven IaC solution"),
            ("3. Project Management", "Agile methodology, communication tools, and development workflow"),
            ("4. Technical Implementation", "Architecture, design patterns, CI/CD integration, and deployment"),
            ("5. Results & Validations", "Implementation results, testing outcomes, and performance metrics"),
            ("6. Future Enhancements", "Current limitations, future work, and conclusion"),
        ]

        self._add_card_grid(slide, cards_data, cols=2, top=Inches(1.5))

    def add_divider_slide(self, title_text):
        """Divider slide with dark teal gradient"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Dark teal gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = TEAL_PRIMARY
        fill.gradient_stops[1].color.rgb = TEAL_DARK

        # Large centered title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(56)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE
        title_frame.word_wrap = True

    def add_bullet_list_slide(self, title, bullets):
        """Content slide with bullet list"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, title)

        # Bullet list with white cards and teal left border
        top = Inches(1.5)
        for bullet in bullets:
            # Parse strong tags
            if "<strong>" in bullet:
                self._add_bullet_card(slide, bullet, top)
            else:
                self._add_bullet_card(slide, bullet, top)
            top += Inches(0.65)

    def add_card_grid_slide(self, title, cards_data, cols=2):
        """Content slide with card grid"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, title)
        self._add_card_grid(slide, cards_data, cols, top=Inches(1.5))

    def add_two_column_slide(self, title, left_bullets, right_content_type, right_data):
        """Two-column slide with text and image/cards"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, title)

        # Left column - bullets
        left_top = Inches(1.5)
        for bullet in left_bullets:
            self._add_bullet_card(slide, bullet, left_top, left=Inches(0.5), width=Inches(4.3))
            left_top += Inches(0.65)

        # Right column - image or content
        if right_content_type == "image":
            self._add_image(slide, right_data, Inches(5.3), Inches(1.5), Inches(4.2), Inches(3.5))
        elif right_content_type == "highlight_box":
            self._add_highlight_box(slide, right_data["title"], right_data["content"],
                                   Inches(5.3), Inches(1.5), Inches(4.2))

    def add_full_image_slide(self, image_name):
        """Full-screen image slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        img_path = os.path.join(self.images_dir, image_name)
        if os.path.exists(img_path):
            # Center image with shadow effect
            self._add_image(slide, image_name, Inches(0.5), Inches(0.4), Inches(9), Inches(4.8))

    def add_tool_grid_slide(self, title, tools):
        """4-column tool grid slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, title)

        # 4-column grid layout - adjusted to fit 6 tools in 2 rows
        col_width = Inches(2.2)
        row_height = Inches(1.8)
        start_left = Inches(0.5)
        start_top = Inches(1.5)

        for idx, tool in enumerate(tools):
            row = idx // 4
            col = idx % 4
            left = start_left + (col * col_width)
            top = start_top + (row * row_height)

            # White card background
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(2.0), Inches(1.6)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(226, 232, 240)
            card.shadow.inherit = False

            # Tool icon/logo
            if tool.get("image"):
                img_path = os.path.join(self.images_dir, tool["image"])
                if os.path.exists(img_path):
                    slide.shapes.add_picture(img_path, left + Inches(0.7), top + Inches(0.2), height=Inches(0.5))
            elif tool.get("placeholder"):
                # Create gradient placeholder for Git
                placeholder = slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    left + Inches(0.7), top + Inches(0.2), Inches(0.5), Inches(0.5)
                )
                placeholder.fill.gradient()
                placeholder.fill.gradient_stops[0].color.rgb = TEAL_PRIMARY
                placeholder.fill.gradient_stops[1].color.rgb = TEAL_LIGHT
                placeholder.line.fill.background()

            # Tool name
            name_box = slide.shapes.add_textbox(left, top + Inches(0.85), Inches(2.0), Inches(0.3))
            name_frame = name_box.text_frame
            name_frame.text = tool["name"]
            name_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            name_frame.paragraphs[0].font.size = Pt(14)
            name_frame.paragraphs[0].font.bold = True
            name_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

            # Description
            desc_box = slide.shapes.add_textbox(left + Inches(0.1), top + Inches(1.15), Inches(1.8), Inches(0.4))
            desc_frame = desc_box.text_frame
            desc_frame.text = tool["description"]
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(9)
            desc_frame.paragraphs[0].font.color.rgb = TEXT_LIGHT
            desc_frame.word_wrap = True

    def add_timeline_slide(self):
        """Timeline slide with 3-column phase grid"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, "Project Timeline & Key Milestones")

        phases = [
            ("Phase 1: Research (2 weeks)", ["OCI API exploration", "Pulumi SDK study", "Requirements gathering", "Architecture design"], RGBColor(59, 130, 246)),
            ("Phase 2: Prototype (2 weeks)", ["Core framework skeleton", "VCN + Instance resources", "Factory pattern POC", "Basic dependency resolver"], RGBColor(139, 92, 246)),
            ("Phase 3: Development (4 weeks)", ["All 7 resource types", "JSON configuration loader", "Dependency resolution logic", "Error handling"], RGBColor(236, 72, 153)),
            ("Phase 4: CI/CD Integration (2 weeks)", ["GitLab pipeline config", "Plan & deploy stages", "Manual approval gates", "Pipeline testing"], RGBColor(16, 185, 129)),
            ("Phase 5: Testing (2 weeks)", ["Integration testing", "Edge case validation", "Performance testing", "Bug fixes"], RGBColor(245, 158, 11)),
            ("Phase 6: Documentation (1 week)", ["User guide creation", "Code documentation", "Architecture diagrams", "Final presentation"], TEAL_PRIMARY),
        ]

        # 3 columns x 2 rows
        col_width = Inches(3.1)
        row_height = Inches(1.8)
        start_left = Inches(0.5)
        start_top = Inches(1.5)

        for idx, (title, items, color) in enumerate(phases):
            row = idx // 3
            col = idx % 3
            left = start_left + (col * col_width)
            top = start_top + (row * row_height)

            # Card with colored top border
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(3.0), Inches(1.6)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(226, 232, 240)

            # Colored top bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, top, Inches(3.0), Inches(0.08)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = color
            top_bar.line.fill.background()

            # Title
            title_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.12), Inches(2.7), Inches(0.25))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(12)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

            # Bullet items
            items_box = slide.shapes.add_textbox(left + Inches(0.25), top + Inches(0.42), Inches(2.6), Inches(1.0))
            items_frame = items_box.text_frame
            for item in items:
                p = items_frame.add_paragraph()
                p.text = "• " + item
                p.font.size = Pt(9)
                p.font.color.rgb = TEXT_GRAY
                p.level = 0

    def add_performance_metrics_slide(self):
        """Performance metrics with large numbers and comparison"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG

        self._add_content_title(slide, "Performance Metrics & Improvements")

        # Three big metric boxes
        metrics = [
            ("95%", "Time Reduction", "From 2-3 hours manual setup\nTo <5 minutes automated", TEAL_PRIMARY, TEAL_DARK),
            ("100%", "Consistency", "Zero configuration drift\nAll envs from same templates", RGBColor(15, 118, 110), RGBColor(20, 184, 166)),
            ("85%", "Error Reduction", "Automated validation\nDependency checking", RED_ACCENT, RGBColor(220, 38, 38)),
        ]

        metric_width = Inches(3.0)
        start_left = Inches(0.5)
        top = Inches(1.5)

        for idx, (number, title, desc, color1, color2) in enumerate(metrics):
            left = start_left + (idx * metric_width)

            # Gradient box
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, top, Inches(2.9), Inches(1.4)
            )
            box.fill.gradient()
            box.fill.gradient_angle = 135
            box.fill.gradient_stops[0].color.rgb = color1
            box.fill.gradient_stops[1].color.rgb = color2
            box.line.fill.background()

            # Big number
            num_box = slide.shapes.add_textbox(left, top + Inches(0.1), Inches(2.9), Inches(0.6))
            num_frame = num_box.text_frame
            num_frame.text = number
            num_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            num_frame.paragraphs[0].font.size = Pt(48)
            num_frame.paragraphs[0].font.bold = True
            num_frame.paragraphs[0].font.color.rgb = WHITE

            # Title
            title_box = slide.shapes.add_textbox(left, top + Inches(0.65), Inches(2.9), Inches(0.25))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            title_frame.paragraphs[0].font.size = Pt(16)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = WHITE

            # Description
            desc_box = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.95), Inches(2.6), Inches(0.4))
            desc_frame = desc_box.text_frame
            desc_frame.text = desc
            desc_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            desc_frame.paragraphs[0].font.size = Pt(10)
            desc_frame.paragraphs[0].font.color.rgb = WHITE
            desc_frame.word_wrap = True

        # Before/After comparison cards
        comparison = [
            ("Before (Manual)", ["Average setup: 2.5 hours", "Error rate: ~15%", "Documentation: Often outdated", "Reproducibility: Low"]),
            ("After (Automated)", ["Average setup: 4 minutes", "Error rate: <2%", "Documentation: Self-documenting JSON", "Reproducibility: 100%"]),
        ]

        comp_top = Inches(3.1)
        for idx, (title, items) in enumerate(comparison):
            left = Inches(0.5) + (idx * Inches(4.7))

            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, comp_top, Inches(4.5), Inches(1.3)
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(226, 232, 240)

            # Top teal bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, comp_top, Inches(4.5), Inches(0.08)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = TEAL_PRIMARY
            top_bar.line.fill.background()

            # Title
            title_box = slide.shapes.add_textbox(left + Inches(0.2), comp_top + Inches(0.12), Inches(4.1), Inches(0.2))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(13)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

            # Items
            items_box = slide.shapes.add_textbox(left + Inches(0.3), comp_top + Inches(0.38), Inches(4.0), Inches(0.85))
            items_frame = items_box.text_frame
            for item in items:
                p = items_frame.add_paragraph()
                p.text = "• " + item
                p.font.size = Pt(10)
                p.font.color.rgb = TEXT_GRAY

    def add_thank_you_slide(self):
        """Final thank you slide with logos"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Dark teal gradient background
        background = slide.background
        fill = background.fill
        fill.gradient()
        fill.gradient_angle = 135
        fill.gradient_stops[0].color.rgb = TEAL_PRIMARY
        fill.gradient_stops[1].color.rgb = TEAL_DARK

        # Thank You title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(1.2), Inches(8), Inches(0.8))
        title_frame = title_box.text_frame
        title_frame.text = "Thank You"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Pt(56)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE

        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(1), Inches(2.1), Inches(8), Inches(0.5))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Questions & Discussion"
        subtitle_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        subtitle_frame.paragraphs[0].font.size = Pt(28)
        subtitle_frame.paragraphs[0].font.color.rgb = WHITE

        # Author info
        author_box = slide.shapes.add_textbox(Inches(1), Inches(3.0), Inches(8), Inches(1.2))
        author_frame = author_box.text_frame
        author_text = "Yassine DBAICHI\nInfrastructure as Code Support in Graal CI\nOracle Labs - GraalVM RISQ Team\nAcademic Year 2024-2025"
        author_frame.text = author_text
        for paragraph in author_frame.paragraphs:
            paragraph.alignment = PP_ALIGN.CENTER
            paragraph.font.size = Pt(18)
            paragraph.font.color.rgb = WHITE
            paragraph.line_spacing = 1.6

        # Logos at bottom
        logo_bottom = Inches(4.6)
        logo1_path = os.path.join(self.images_dir, "logo_ehtp.jpg")
        logo2_path = os.path.join(self.images_dir, "logo_graalvm.png")

        if os.path.exists(logo1_path):
            slide.shapes.add_picture(logo1_path, Inches(3.5), logo_bottom, height=Inches(0.65))
        if os.path.exists(logo2_path):
            slide.shapes.add_picture(logo2_path, Inches(5.8), logo_bottom, height=Inches(0.65))

    # Helper methods
    def _add_content_title(self, slide, title_text):
        """Add title with teal color and red underline"""
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(9), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.text = title_text
        title_frame.paragraphs[0].font.size = Pt(32)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

        # Red underline
        underline = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(0.5), Inches(0.95), Inches(9), Inches(0.08)
        )
        underline.fill.solid()
        underline.fill.fore_color.rgb = RED_ACCENT
        underline.line.fill.background()

    def _add_bullet_card(self, slide, text, top, left=Inches(0.5), width=Inches(9)):
        """Add white card with teal left border for bullet point"""
        card_height = Inches(0.55)

        # White card background
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, card_height
        )
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        # Teal left border
        border = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            left, top, Inches(0.08), card_height
        )
        border.fill.solid()
        border.fill.fore_color.rgb = TEAL_PRIMARY
        border.line.fill.background()

        # Text
        text_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.05), width - Inches(0.3), card_height - Inches(0.1))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True
        text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        # Parse bold text
        if "<strong>" in text:
            parts = text.split("<strong>")
            p = text_frame.paragraphs[0]
            if parts[0]:
                p.text = parts[0]
                p.font.size = Pt(14)
                p.font.color.rgb = TEXT_GRAY

            for part in parts[1:]:
                bold_text, normal_text = part.split("</strong>", 1)

                run = p.add_run()
                run.text = bold_text
                run.font.size = Pt(14)
                run.font.bold = True
                run.font.color.rgb = TEXT_GRAY

                if normal_text:
                    run = p.add_run()
                    run.text = normal_text
                    run.font.size = Pt(14)
                    run.font.color.rgb = TEXT_GRAY
        else:
            text_frame.text = text
            text_frame.paragraphs[0].font.size = Pt(14)
            text_frame.paragraphs[0].font.color.rgb = TEXT_GRAY

    def _add_card_grid(self, slide, cards_data, cols, top):
        """Add grid of cards"""
        card_width = Inches((9 / cols) - 0.2)
        card_height = Inches(1.2)
        rows = (len(cards_data) + cols - 1) // cols

        for idx, (title, content) in enumerate(cards_data):
            row = idx // cols
            col = idx % cols
            left = Inches(0.5) + (col * (card_width + Inches(0.2)))
            card_top = top + (row * (card_height + Inches(0.15)))

            # White card with teal top border
            card = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left, card_top, card_width, card_height
            )
            card.fill.solid()
            card.fill.fore_color.rgb = WHITE
            card.line.color.rgb = RGBColor(226, 232, 240)

            # Teal top bar
            top_bar = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                left, card_top, card_width, Inches(0.08)
            )
            top_bar.fill.solid()
            top_bar.fill.fore_color.rgb = TEAL_PRIMARY
            top_bar.line.fill.background()

            # Card title
            title_box = slide.shapes.add_textbox(left + Inches(0.15), card_top + Inches(0.12), card_width - Inches(0.3), Inches(0.25))
            title_frame = title_box.text_frame
            title_frame.text = title
            title_frame.paragraphs[0].font.size = Pt(13)
            title_frame.paragraphs[0].font.bold = True
            title_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY
            title_frame.word_wrap = True

            # Card content
            content_box = slide.shapes.add_textbox(left + Inches(0.15), card_top + Inches(0.42), card_width - Inches(0.3), Inches(0.7))
            content_frame = content_box.text_frame
            content_frame.text = content
            content_frame.paragraphs[0].font.size = Pt(11)
            content_frame.paragraphs[0].font.color.rgb = TEXT_GRAY
            content_frame.word_wrap = True

    def _add_highlight_box(self, slide, title, content, left, top, width):
        """Add gradient highlight box"""
        box_height = Inches(1.5)

        # Gradient box
        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            left, top, width, box_height
        )
        box.fill.gradient()
        box.fill.gradient_angle = 135
        box.fill.gradient_stops[0].color.rgb = TEAL_PRIMARY
        box.fill.gradient_stops[1].color.rgb = TEAL_LIGHT
        box.line.fill.background()

        # Title
        title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.3))
        title_frame = title_box.text_frame
        title_frame.text = title
        title_frame.paragraphs[0].font.size = Pt(18)
        title_frame.paragraphs[0].font.bold = True
        title_frame.paragraphs[0].font.color.rgb = WHITE

        # Content
        content_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.5), width - Inches(0.4), Inches(0.9))
        content_frame = content_box.text_frame
        content_frame.text = content
        content_frame.paragraphs[0].font.size = Pt(12)
        content_frame.paragraphs[0].font.color.rgb = WHITE
        content_frame.word_wrap = True

    def _add_image(self, slide, image_name, left, top, width, height):
        """Add image to slide"""
        img_path = os.path.join(self.images_dir, image_name)
        if os.path.exists(img_path):
            try:
                slide.shapes.add_picture(img_path, left, top, width=width, height=height)
            except:
                # If image can't fit exact dimensions, let it scale
                slide.shapes.add_picture(img_path, left, top, width=width)

    def save(self, output_path):
        """Save presentation to file"""
        self.prs.save(output_path)
        return output_path


def main():
    images_dir = "/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/images"
    output_path = "/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation_final_perfect.pptx"

    print("Generating pixel-perfect PowerPoint presentation...")
    print(f"Images directory: {images_dir}")

    builder = PresentationBuilder(images_dir)

    # Slide 1: Title
    print("  [1/58] Title slide...")
    builder.add_title_slide()

    # Slide 2: Agenda
    print("  [2/58] Agenda...")
    builder.add_agenda_slide()

    # SECTION 1: COMPANY & TEAM CONTEXT
    print("  [3/58] Divider: Company & Team Context...")
    builder.add_divider_slide("Company & Team Context")

    print("  [4/58] Oracle Corporation...")
    builder.add_bullet_list_slide("Oracle Corporation", [
        "<strong>Global Technology Leader:</strong> Multinational computer technology corporation specializing in database software and cloud solutions",
        "<strong>Founded:</strong> 1977 by Larry Ellison, Bob Miner, and Ed Oates",
        "<strong>Headquarters:</strong> Austin, Texas, USA (Oracle Headquarters)",
        "<strong>Market Position:</strong> One of the largest enterprise software companies worldwide",
        "<strong>Core Products:</strong> Oracle Database, Oracle Cloud Infrastructure (OCI), Enterprise Applications",
        "<strong>Revenue:</strong> $50+ billion annually with 430,000+ employees globally"
    ])

    print("  [5/58] Oracle Labs...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "Oracle Labs - Research & Innovation")

    # Left column bullets
    left_bullets = [
        "<strong>Mission:</strong> Conduct research to advance the state of the art in computer science and systems",
        "<strong>Focus Areas:</strong> Programming languages, Virtual machines, Database systems, Cloud infrastructure",
        "<strong>Key Projects:</strong> GraalVM, Truffle Framework, SubstrateVM"
    ]
    top = Inches(1.5)
    for bullet in left_bullets:
        builder._add_bullet_card(slide, bullet, top, Inches(0.5), Inches(4.3))
        top += Inches(0.65)

    # Right column - highlight box + card
    builder._add_highlight_box(slide, "Research Excellence",
                               "Oracle Labs bridges cutting-edge academic research with real-world enterprise applications\n\nOutput: Open-source projects, academic publications, and production-ready technologies",
                               Inches(5.3), Inches(1.5), Inches(4.2))

    # Small card below
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(3.3), Inches(4.2), Inches(0.8))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(3.3), Inches(4.2), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(5.45), Inches(3.42), Inches(3.9), Inches(0.2))
    title_box.text_frame.text = "Collaboration"
    title_box.text_frame.paragraphs[0].font.size = Pt(13)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    content_box = slide.shapes.add_textbox(Inches(5.45), Inches(3.65), Inches(3.9), Inches(0.4))
    content_box.text_frame.text = "Strong partnerships with universities and research institutions worldwide"
    content_box.text_frame.paragraphs[0].font.size = Pt(11)
    content_box.text_frame.paragraphs[0].font.color.rgb = TEXT_GRAY
    content_box.text_frame.word_wrap = True

    print("  [6/58] GraalVM Overview...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "GraalVM - High-Performance Polyglot Runtime")

    # Left bullets
    left_bullets = [
        "<strong>What is GraalVM?</strong> A universal virtual machine for running applications written in JavaScript, Python, Ruby, R, JVM languages, and more",
        "<strong>Key Features:</strong> Polyglot capabilities, Native Image compilation, Superior performance optimization, Reduced memory footprint",
        "<strong>Use Cases:</strong> Microservices, serverless computing, containerized applications"
    ]
    top = Inches(1.5)
    for bullet in left_bullets:
        builder._add_bullet_card(slide, bullet, top, Inches(0.5), Inches(4.3))
        top += Inches(0.65)

    # Right image
    builder._add_image(slide, "logo_graalvm.png", Inches(5.3), Inches(1.8), Inches(4.2), Inches(3.0))

    print("  [7/58] GraalVM RISQ Team...")
    builder.add_bullet_list_slide("GraalVM RISQ Team", [
        "<strong>Team Name:</strong> RISQ (Release Infrastructure, Systems & Quality)",
        "<strong>Mission:</strong> Ensure reliability, scalability, and quality of GraalVM infrastructure and development workflows",
        "<strong>Responsibilities:</strong> CI/CD pipeline management, Infrastructure automation, Testing infrastructure, Developer tooling improvements",
        "<strong>Tech Stack:</strong> GitLab CI/CD, Oracle Cloud Infrastructure, Python, Bash scripting"
    ])

    # Add highlight box manually
    slide = builder.prs.slides[builder.prs.slides.index(builder.prs.slides[-1])]
    builder._add_highlight_box(slide, "",
                               "The RISQ team acts as the backbone supporting GraalVM development, ensuring smooth operations for hundreds of developers",
                               Inches(0.5), Inches(4.4), Inches(9))

    print("  [8/58] Project Context & My Role...")
    builder.add_card_grid_slide("Project Context & Internship Focus", [
        ("Team Challenge", "GraalVM requires frequent infrastructure provisioning for testing, CI/CD, and development environments across multiple cloud regions"),
        ("Manual Bottleneck", "Infrastructure management was manual, error-prone, and time-consuming - limiting team agility"),
        ("Project Goal", "Design and implement an automated Infrastructure as Code (IaC) solution integrated with Graal CI pipeline"),
        ("My Role", "Lead developer responsible for architecture design, implementation, and GitLab CI/CD integration")
    ], cols=2)

    # SECTION 2: PROBLEM & SOLUTION
    print("  [9/58] Divider: Problem & Solution Overview...")
    builder.add_divider_slide("Problem & Solution Overview")

    print("  [10/58] Current Challenges...")
    builder.add_bullet_list_slide("Current Infrastructure Management Challenges", [
        "<strong>Manual Console Operations:</strong> Team members had to manually create resources through OCI web console - time-consuming and repetitive",
        "<strong>Lack of Version Control:</strong> No tracking of infrastructure changes, making rollbacks and audits difficult",
        "<strong>Configuration Drift:</strong> Resources created manually often diverged from documented standards",
        "<strong>Scalability Issues:</strong> Provisioning multiple similar environments was labor-intensive",
        "<strong>Knowledge Silos:</strong> Infrastructure knowledge concentrated in few team members",
        "<strong>No CI/CD Integration:</strong> Infrastructure provisioning was separate from development workflow"
    ])

    slide = builder.prs.slides[builder.prs.slides.index(builder.prs.slides[-1])]
    builder._add_highlight_box(slide, "",
                               "Result: Slow deployment cycles, increased error rates, and reduced team productivity",
                               Inches(0.5), Inches(4.8), Inches(9))

    print("  [11/58] Functional Requirements...")
    builder.add_bullet_list_slide("Functional Requirements", [
        "<strong>FR1 - Declarative Configuration:</strong> Users define infrastructure in human-readable format (JSON) without writing Python/Terraform code",
        "<strong>FR2 - Resource Support:</strong> Support for essential OCI resources (VCN, Subnets, Compute Instances, Storage Buckets, Security Lists, Routing)",
        "<strong>FR3 - Dependency Management:</strong> Automatic resolution of resource dependencies (e.g., subnet requires VCN, instance requires subnet)",
        "<strong>FR4 - CI/CD Integration:</strong> Seamless integration with GitLab CI/CD pipeline (plan, preview, apply stages)",
        "<strong>FR5 - State Management:</strong> Track infrastructure state to enable updates and deletions",
        "<strong>FR6 - Multi-Environment:</strong> Support for development, staging, and production configurations"
    ])

    print("  [12/58] Non-Functional Requirements...")
    builder.add_card_grid_slide("Non-Functional Requirements", [
        ("Performance", "• Provision infrastructure in <5 minutes\n• Support concurrent deployments\n• Minimal API calls to OCI"),
        ("Security", "• Credentials stored securely (GitLab secrets)\n• No hardcoded sensitive data\n• Audit trail via Git commits"),
        ("Extensibility", "• Plugin architecture for new resources\n• Support for future cloud providers\n• Custom validation rules"),
    ], cols=2)

    print("  [13/58] Technology Comparison...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "IaC Technology Selection: Terraform vs Pulumi")

    # Two large comparison cards
    # Terraform card (left)
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.4), Inches(3.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = RED_ACCENT
    card1.line.width = Pt(2)

    top_bar1 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(4.4), Inches(0.1))
    top_bar1.fill.solid()
    top_bar1.fill.fore_color.rgb = RED_ACCENT
    top_bar1.line.fill.background()

    title1 = slide.shapes.add_textbox(Inches(0.65), Inches(1.65), Inches(4.1), Inches(0.25))
    title1.text_frame.text = "Terraform"
    title1.text_frame.paragraphs[0].font.size = Pt(16)
    title1.text_frame.paragraphs[0].font.bold = True
    title1.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    content1 = slide.shapes.add_textbox(Inches(0.65), Inches(1.95), Inches(4.1), Inches(2.75))
    tf = content1.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "Pros:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_GRAY

    pros = ["Industry standard with mature ecosystem", "HCL declarative language", "Large provider ecosystem", "Strong community support"]
    for pro in pros:
        p = tf.add_paragraph()
        p.text = "• " + pro
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY
        p.level = 0

    p = tf.add_paragraph()
    p.text = "\nCons:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_GRAY

    cons = ["Custom DSL (HCL) requires learning", "Limited programming logic & abstraction", "State management complexity", "Hard to build dynamic configurations"]
    for con in cons:
        p = tf.add_paragraph()
        p.text = "• " + con
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY
        p.level = 0

    # Pulumi card (right) - selected
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), Inches(1.5), Inches(4.4), Inches(3.3))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = TEAL_LIGHT
    card2.line.width = Pt(3)

    top_bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), Inches(1.5), Inches(4.4), Inches(0.1))
    top_bar2.fill.solid()
    top_bar2.fill.fore_color.rgb = TEAL_LIGHT
    top_bar2.line.fill.background()

    title2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.65), Inches(4.1), Inches(0.25))
    title2.text_frame.text = "Pulumi ✓ (Selected)"
    title2.text_frame.paragraphs[0].font.size = Pt(16)
    title2.text_frame.paragraphs[0].font.bold = True
    title2.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    content2 = slide.shapes.add_textbox(Inches(5.25), Inches(1.95), Inches(4.1), Inches(2.75))
    tf2 = content2.text_frame
    tf2.word_wrap = True

    p = tf2.paragraphs[0]
    p.text = "Pros:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = TEXT_GRAY

    pros2 = ["Use real programming languages (Python)", "Full programmatic control & abstraction", "Built-in state management", "Excellent OCI support", "Enables framework development on top", "Team Python expertise leverage"]
    for pro in pros2:
        p = tf2.add_paragraph()
        p.text = "• " + pro
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY

    p = tf2.add_paragraph()
    p.text = "\nWhy chosen: Perfect fit for building JSON-driven abstraction layer. Pulumi's programmatic approach allows us to create a framework that hides IaC complexity from users."
    p.font.size = Pt(9)
    p.font.color.rgb = TEAL_PRIMARY
    p.font.italic = True

    print("  [14/58] How does Pulumi work...")
    builder.add_card_grid_slide("How does Pulumi work?", [
        ("1. Write Code", "Define infrastructure using Python - loops, functions, conditionals available"),
        ("2. Preview", "Run 'pulumi preview' to see what will change before applying"),
        ("3. Deploy", "Pulumi SDK translates Python to OCI API calls and provisions resources"),
        ("4. State Management", "Pulumi tracks all resources in state file (OCI Object Storage)")
    ], cols=2)

    # Add bullets below
    slide = builder.prs.slides[builder.prs.slides.index(builder.prs.slides[-1])]
    bullets_below = [
        "<strong>Python-First:</strong> Leverage existing team Python expertise",
        "<strong>OCI Provider:</strong> First-class Oracle Cloud Infrastructure support",
        "<strong>Framework Foundation:</strong> Enables building JSON-driven abstraction layer on top"
    ]
    top = Inches(3.8)
    for bullet in bullets_below:
        builder._add_bullet_card(slide, bullet, top, Inches(0.5), Inches(9))
        top += Inches(0.55)

    print("  [15/58] Framework Solution Pillars...")
    builder.add_full_image_slide("the framework pillars , features.png")

    print("  [16/58] Layer Architecture...")
    builder.add_full_image_slide("arch_layer architecture_system.png")

    # SECTION 3: PROJECT MANAGEMENT
    print("  [17/58] Divider: Project Management...")
    builder.add_divider_slide("Project Management")

    print("  [18/58] Management Approach...")
    builder.add_bullet_list_slide("Project Management Approach", [
        "<strong>Methodology:</strong> Kanban Agile - visual workflow management with continuous delivery",
        "<strong>Project Tracker:</strong> Task board for tracking work items (To Do, In Progress, Done)",
        "<strong>Regular Meetings:</strong> Weekly sync meetings with mentor to review progress and blockers",
        "<strong>Continuous Flow:</strong> Work items pulled as capacity allows, no fixed sprint boundaries",
        "<strong>Focus on Delivery:</strong> Prioritize completing tasks over starting new ones"
    ])

    print("  [19/58] Development Workflow...")
    builder.add_full_image_slide("mgmt_task_workflow.png")

    print("  [20/58] Communication Tools...")
    builder.add_tool_grid_slide("Communication & Collaboration Tools", [
        {"name": "Slack", "image": "tool_slack.png", "description": "Daily communication, quick questions"},
        {"name": "Zoom", "image": "tool_zoom.png", "description": "Weekly sync meetings, discussions"},
        {"name": "Jira", "image": "tool_jira.png", "description": "Task tracking, sprint planning"},
        {"name": "Outlook", "image": "tool_outlook.png", "description": "Meeting scheduling, formal comms"},
        {"name": "Confluence", "image": "confluence logo.png", "description": "Documentation, knowledge base"},
        {"name": "GitLab", "placeholder": True, "description": "Code repository, CI/CD, MRs"}
    ])

    print("  [21/58] Mentor Sessions...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "Mentor Check-ins & Team Meetings")

    # Left - image
    builder._add_image(slide, "mgmt_mentor_session.png", Inches(0.5), Inches(1.5), Inches(4.3), Inches(3.0))

    # Right - cards
    card1 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(1.5), Inches(4.2), Inches(1.3))
    card1.fill.solid()
    card1.fill.fore_color.rgb = WHITE
    card1.line.color.rgb = RGBColor(226, 232, 240)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(1.5), Inches(4.2), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(5.45), Inches(1.62), Inches(3.9), Inches(0.2))
    title_box.text_frame.text = "Weekly Sync Meetings"
    title_box.text_frame.paragraphs[0].font.size = Pt(13)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    content_box = slide.shapes.add_textbox(Inches(5.55), Inches(1.88), Inches(3.8), Inches(0.85))
    cf = content_box.text_frame
    cf.word_wrap = True
    items = ["Frequency: Every Tuesday 11:00 AM", "Duration: 45-60 minutes", "Agenda: Progress review, blockers discussion, next steps planning", "Participants: Intern + mentor (Hamza Ghaissi)"]
    for item in items:
        p = cf.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(10)
        p.font.color.rgb = TEXT_GRAY

    # Card 2
    card2 = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.3), Inches(3.0), Inches(4.2), Inches(1.0))
    card2.fill.solid()
    card2.fill.fore_color.rgb = WHITE
    card2.line.color.rgb = RGBColor(226, 232, 240)

    top_bar2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.3), Inches(3.0), Inches(4.2), Inches(0.08))
    top_bar2.fill.solid()
    top_bar2.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar2.line.fill.background()

    title_box2 = slide.shapes.add_textbox(Inches(5.45), Inches(3.12), Inches(3.9), Inches(0.2))
    title_box2.text_frame.text = "Ad-hoc Sessions"
    title_box2.text_frame.paragraphs[0].font.size = Pt(13)
    title_box2.text_frame.paragraphs[0].font.bold = True
    title_box2.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    content_box2 = slide.shapes.add_textbox(Inches(5.55), Inches(3.38), Inches(3.8), Inches(0.55))
    cf2 = content_box2.text_frame
    cf2.word_wrap = True
    items2 = ["Technical architecture discussions", "Code reviews and pair programming", "Debugging complex issues", "Design decision consultations"]
    for item in items2:
        p = cf2.add_paragraph()
        p.text = "• " + item
        p.font.size = Pt(9)
        p.font.color.rgb = TEXT_GRAY

    # Highlight box
    builder._add_highlight_box(slide, "",
                               "Communication Style: Open-door policy via Slack - quick responses and collaborative problem-solving",
                               Inches(5.3), Inches(4.15), Inches(4.2))

    print("  [22/58] Timeline & Milestones...")
    builder.add_timeline_slide()

    # SECTION 4: TECHNICAL IMPLEMENTATION
    print("  [23/58] Divider: Technical Implementation...")
    builder.add_divider_slide("Technical Implementation")

    print("  [24/58] Use Case Diagram...")
    builder.add_full_image_slide("design_usecase_diagram.png")

    print("  [25/58] Project Structure...")
    builder.add_full_image_slide("project folders .png")

    print("  [26/58] JSON Configuration Example...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "JSON Configuration Example")

    # Code card
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(2.0))
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = RGBColor(226, 232, 240)

    top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(1.5), Inches(9), Inches(0.08))
    top_bar.fill.solid()
    top_bar.fill.fore_color.rgb = TEAL_PRIMARY
    top_bar.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.62), Inches(8.7), Inches(0.2))
    title_box.text_frame.text = "infrastructure.json"
    title_box.text_frame.paragraphs[0].font.size = Pt(13)
    title_box.text_frame.paragraphs[0].font.bold = True
    title_box.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

    # JSON code
    code_box = slide.shapes.add_textbox(Inches(0.8), Inches(1.9), Inches(8.4), Inches(1.5))
    code_box.fill.solid()
    code_box.fill.fore_color.rgb = RGBColor(241, 245, 249)
    cf = code_box.text_frame
    cf.word_wrap = False
    json_code = '''{
  "name": "production-vcn",
  "cidrBlock": "10.0.0.0/16",
  "internetGateways": [...],
  "securityLists": [...],
  "routeTables": [...],
  "subnets": [...]
}'''
    cf.text = json_code
    cf.paragraphs[0].font.name = "Courier New"
    cf.paragraphs[0].font.size = Pt(13)
    cf.paragraphs[0].font.color.rgb = TEXT_GRAY

    # Bullets below
    bullets = [
        "<strong>Name-Based References:</strong> Resources reference each other by human-readable names",
        "<strong>Automatic Resolution:</strong> Framework resolves names to OCIDs automatically"
    ]
    top = Inches(3.7)
    for bullet in bullets:
        builder._add_bullet_card(slide, bullet, top)
        top += Inches(0.55)

    print("  [27/58] Core Framework Components...")
    builder.add_card_grid_slide("Core Framework Components", [
        ("BaseResource (Abstract)", "File: core/base_resource.py (21 lines)\nPurpose: Abstract base class defining resource interface\nMethods: create() - Factory method, build() - Abstract (must implement)"),
        ("ResourceFactory", "File: core/resource_factory.py (38 lines)\nPurpose: Factory pattern for dynamic instantiation\nRegistry: Maps type strings → Python classes\n7 resource types registered"),
        ("DependencyResolver", "File: core/dependency_resolver.py (14 lines)\nPurpose: Resource reference manager\nMethods: register_resource(name, resource), get_resource_id(name)"),
        ("ConfigLoader", "File: config/config_loader.py (25 lines)\nPurpose: Parse infrastructure.json\nMethods: get_vcns_config(), get_instances_config(), get_image_config()")
    ], cols=2)

    print("  [28/58] Class Diagram...")
    builder.add_full_image_slide("class diagram .png")

    print("  [29/58] Factory Pattern...")
    builder.add_full_image_slide("arch_factory_pattern.png")

    print("  [30/58] Dependency Resolver...")
    builder.add_full_image_slide("arch_dependency_resolution.png")

    print("  [31/58] Implemented Resources...")
    builder.add_card_grid_slide("Implemented Resource Types", [
        ("1. VCN (Virtual Cloud Network)", "Purpose: Foundation network container"),
        ("2. Internet Gateway", "Purpose: Enable internet connectivity"),
        ("3. Security List", "Purpose: Firewall rules (ingress/egress)"),
        ("4. Route Table", "Purpose: Network routing configuration"),
        ("5. Subnet", "Purpose: Network segment within VCN"),
        ("6. Compute Instance", "Purpose: Virtual machine"),
        ("7. Object Storage Bucket", "Purpose: S3-compatible object storage")
    ], cols=2)

    print("  [32/58] Workflow Design...")
    builder.add_full_image_slide("workflow_design.png")

    print("  [33/58] GitLab Pipeline Stages...")
    builder.add_full_image_slide("cicd_gitlab_pipeline plan and deploy stages.png")

    print("  [34/58] Pull Request Workflow...")
    builder.add_full_image_slide("workflow_pull_request.png")

    # SECTION 5: RESULTS & VALIDATIONS
    print("  [35/58] Divider: Results & Validations...")
    builder.add_divider_slide("Results & Validations")

    print("  [36/58] Implementation Results...")
    builder.add_card_grid_slide("Implementation Results", [
        ("✓ Deliverables Completed", "• Python framework (800 LOC)\n• 7 OCI resource types\n• JSON configuration system\n• GitLab CI/CD pipeline\n• Documentation & diagrams"),
        ("✓ Requirements Met", "• Declarative JSON config\n• Dependency resolution\n• CI/CD integration\n• Multi-environment support\n• State management"),
        ("✓ Production Ready", "• Error handling implemented\n• Tested on dev & staging\n• Documentation complete\n• Team training delivered\n• Now in active use by RISQ team")
    ], cols=2)

    print("  [37/58] Deployment Success...")
    builder.add_full_image_slide("result_instance_deployed.png")

    print("  [38/58] Testing Results...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "Testing & Validation Results")

    # Two column layout with 4 cards total
    cards_left = [
        ("Unit Testing", "Scope: Individual resource classes\n• Factory pattern registration\n• Dependency resolver lookup\n• Config loader parsing\n\n✓ All components working correctly"),
        ("Integration Testing", "Test Cases:\n• Single VCN + 1 instance\n• 2 VCNs + 2 instances\n• Complex network (IGW + RT + SL)\n• Bucket creation\n\n✓ End-to-end provisioning successful")
    ]

    cards_right = [
        ("CI/CD Pipeline Testing", "Validated:\n• Plan stage (preview changes)\n• Deploy stage (apply changes)\n• Manual approval gate\n• State persistence\n\n✓ Pipeline stable over 20+ runs"),
        ("Edge Cases Tested", "• Missing dependency (error caught)\n• Invalid CIDR block (validation)\n• Duplicate resource names (detected)\n• Image not compatible with shape")
    ]

    # Left column
    for idx, (title, content) in enumerate(cards_left):
        top = Inches(1.5) + (idx * Inches(1.6))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), top, Inches(4.4), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), top, Inches(4.4), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = TEAL_PRIMARY
        top_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(0.65), top + Inches(0.12), Inches(4.1), Inches(0.2))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(13)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

        content_box = slide.shapes.add_textbox(Inches(0.65), top + Inches(0.38), Inches(4.1), Inches(0.95))
        content_box.text_frame.text = content
        content_box.text_frame.paragraphs[0].font.size = Pt(9)
        content_box.text_frame.paragraphs[0].font.color.rgb = TEXT_GRAY
        content_box.text_frame.word_wrap = True

    # Right column
    for idx, (title, content) in enumerate(cards_right):
        top = Inches(1.5) + (idx * Inches(1.6))
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(5.1), top, Inches(4.4), Inches(1.4))
        card.fill.solid()
        card.fill.fore_color.rgb = WHITE
        card.line.color.rgb = RGBColor(226, 232, 240)

        top_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.1), top, Inches(4.4), Inches(0.08))
        top_bar.fill.solid()
        top_bar.fill.fore_color.rgb = TEAL_PRIMARY
        top_bar.line.fill.background()

        title_box = slide.shapes.add_textbox(Inches(5.25), top + Inches(0.12), Inches(4.1), Inches(0.2))
        title_box.text_frame.text = title
        title_box.text_frame.paragraphs[0].font.size = Pt(13)
        title_box.text_frame.paragraphs[0].font.bold = True
        title_box.text_frame.paragraphs[0].font.color.rgb = TEAL_PRIMARY

        content_box = slide.shapes.add_textbox(Inches(5.25), top + Inches(0.38), Inches(4.1), Inches(0.95))
        content_box.text_frame.text = content
        content_box.text_frame.paragraphs[0].font.size = Pt(9)
        content_box.text_frame.paragraphs[0].font.color.rgb = TEXT_GRAY
        content_box.text_frame.word_wrap = True

    print("  [39/58] Performance Metrics...")
    builder.add_performance_metrics_slide()

    print("  [40/58] Business Impact...")
    builder.add_bullet_list_slide("Business Impact & Team Benefits", [
        "<strong>Developer Productivity:</strong> Team members can now focus on GraalVM development instead of infrastructure management - estimated 10+ hours saved per week across the team",
        "<strong>Faster Testing Cycles:</strong> Spin up test environments in minutes instead of hours - enables rapid experimentation and bug fixes",
        "<strong>Onboarding Acceleration:</strong> New team members can provision infrastructure on day one without specialized training - reduces onboarding time from days to hours",
        "<strong>Cost Optimization:</strong> All infrastructure tracked in Git - easier to identify and delete unused resources, preventing budget overruns"
    ])

    # SECTION 6: FUTURE ENHANCEMENTS
    print("  [41/58] Divider: Future Enhancements...")
    builder.add_divider_slide("Future Enhancements")

    print("  [42/58] Future Improvements...")
    builder.add_card_grid_slide("Future Improvements", [
        ("More Resource Coverage", "Expand framework to support additional OCI resources like Load Balancers, Databases, and File Storage"),
        ("Schema Validation", "Implement JSON schema validation to catch configuration errors before deployment"),
        ("Role Identification", "Define clear roles within the team on how to use the framework and establish best practices")
    ], cols=2)

    print("  [43/58] Conclusion...")
    slide = builder.prs.slides.add_slide(builder.prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = LIGHT_GRAY_BG
    builder._add_content_title(slide, "Conclusion")

    builder._add_highlight_box(slide, "",
                               "Successfully delivered a production-ready IaC framework that transforms infrastructure management for the GraalVM RISQ team",
                               Inches(1.5), Inches(2.5), Inches(7))

    print("  [44/58] Thank You...")
    builder.add_thank_you_slide()

    # Save presentation
    print("\nSaving presentation...")
    output = builder.save(output_path)
    print(f"✓ Presentation saved: {output}")
    print(f"\nTotal slides: {len(builder.prs.slides)}")

    return output


if __name__ == "__main__":
    output_file = main()
    print(f"\n{'='*60}")
    print(f"SUCCESS! Pixel-perfect PowerPoint generated:")
    print(f"{output_file}")
    print(f"{'='*60}")
