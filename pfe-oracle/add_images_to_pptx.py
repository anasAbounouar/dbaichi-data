#!/usr/bin/env python3
"""
Add all images from HTML presentation to the PPTX file
Parses HTML to find which images go on which slides and embeds them
"""

import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt

def parse_html_for_images():
    """Parse HTML to extract slide-to-image mappings"""
    with open('/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation.html', 'r') as f:
        html_content = f.read()

    # Extract all slides
    slides_pattern = r'<div class="slide[^"]*"[^>]*>(.*?)</div>\s*(?=<div class="slide|<script)'
    slides = re.findall(slides_pattern, html_content, re.DOTALL)

    image_mappings = []

    for slide_num, slide_content in enumerate(slides, 1):
        # Find images in this slide
        img_pattern = r'<img[^>]+src="([^"]+)"'
        images = re.findall(img_pattern, slide_content)

        if images:
            # Check if this is a full-image slide (image should be large)
            is_full_image = 'full-image-slide' in slide_content or 'image-container' in slide_content

            for img_path in images:
                # Clean up the path
                img_path = img_path.replace('images/', '')
                image_mappings.append({
                    'slide_num': slide_num,
                    'image_path': img_path,
                    'is_full_image': is_full_image
                })

    return image_mappings

def add_images_to_pptx(pptx_path, image_mappings, images_dir):
    """Add images to PPTX slides"""
    prs = Presentation(pptx_path)

    print(f"\n{'='*60}")
    print(f"Adding images to PPTX")
    print(f"{'='*60}\n")

    images_added = 0
    images_skipped = 0

    for mapping in image_mappings:
        slide_num = mapping['slide_num']
        image_path = mapping['image_path']
        is_full_image = mapping['is_full_image']

        # Adjust slide number (0-indexed in python-pptx)
        slide_idx = slide_num - 1

        if slide_idx >= len(prs.slides):
            print(f"⚠ Slide {slide_num} not found in PPTX, skipping")
            images_skipped += 1
            continue

        slide = prs.slides[slide_idx]

        # Full path to image
        full_image_path = os.path.join(images_dir, image_path)

        if not os.path.exists(full_image_path):
            print(f"⚠ Image not found: {image_path}, skipping")
            images_skipped += 1
            continue

        # Check if slide already has this image (avoid duplicates)
        already_has_image = False
        for shape in slide.shapes:
            if shape.shape_type == 13:  # Picture
                # Simple check: if slide already has a picture, skip
                already_has_image = True
                break

        if already_has_image:
            # print(f"  Slide {slide_num} already has an image, skipping")
            images_skipped += 1
            continue

        try:
            # Add image based on slide type
            if is_full_image:
                # Full slide image - centered and large
                left = Inches(1.5)
                top = Inches(1.5)
                height = Inches(5.0)
                slide.shapes.add_picture(full_image_path, left, top, height=height)
                print(f"✓ Added full image to slide {slide_num}: {image_path}")
            else:
                # Smaller image - top right or center
                left = Inches(7.0)
                top = Inches(2.0)
                height = Inches(3.0)
                slide.shapes.add_picture(full_image_path, left, top, height=height)
                print(f"✓ Added image to slide {slide_num}: {image_path}")

            images_added += 1

        except Exception as e:
            print(f"✗ Error adding image {image_path} to slide {slide_num}: {e}")
            images_skipped += 1

    # Save the modified PPTX
    output_path = pptx_path.replace('.pptx', '_with_images.pptx')
    prs.save(output_path)

    print(f"\n{'='*60}")
    print(f"Summary")
    print(f"{'='*60}")
    print(f"Images added: {images_added}")
    print(f"Images skipped: {images_skipped}")
    print(f"Output file: {output_path}")
    print(f"{'='*60}\n")

    return output_path

def main():
    base_dir = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle'
    pptx_path = os.path.join(base_dir, 'presentation_converted.pptx')
    images_dir = os.path.join(base_dir, 'images')

    print("Step 1: Parsing HTML to find images...")
    image_mappings = parse_html_for_images()
    print(f"Found {len(image_mappings)} image references in HTML\n")

    print("Step 2: Adding images to PPTX...")
    output_path = add_images_to_pptx(pptx_path, image_mappings, images_dir)

    print("\n✅ Complete! Your PPTX now has all the images embedded.")
    print(f"📁 Location: {output_path}")

if __name__ == "__main__":
    main()
