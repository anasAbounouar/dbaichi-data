#!/usr/bin/env python3
"""
Convert HTML presentation to PowerPoint
"""
import re
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from lxml import html

# Initialize presentation
prs = Presentation()
prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
prs.slide_height = Inches(7.5)

# Read HTML file
html_file = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation.html'
with open(html_file, 'r', encoding='utf-8') as f:
    content = f.read()

# Parse HTML
tree = html.fromstring(content)

# Helper function to clean text
def clean_text(text):
    if text is None:
        return ""
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text.strip())
    return text

# Helper function to add title slide
def add_title_slide(title, subtitle=""):
    slide_layout = prs.slide_layouts[0]  # Title slide layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    if subtitle and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle

    return slide

# Helper function to add content slide
def add_content_slide(title, body_text=""):
    slide_layout = prs.slide_layouts[1]  # Title and content layout
    slide = prs.slides.add_slide(slide_layout)

    title_shape = slide.shapes.title
    title_shape.text = title

    return slide

# Helper function to add image to slide
def add_image_to_slide(slide, img_path, left, top, width=None, height=None):
    try:
        if os.path.exists(img_path):
            if width and height:
                slide.shapes.add_picture(img_path, left, top, width, height)
            elif width:
                slide.shapes.add_picture(img_path, left, top, width=width)
            else:
                slide.shapes.add_picture(img_path, left, top)
            return True
    except Exception as e:
        print(f"Error adding image {img_path}: {e}")
    return False

# Helper function to add text box
def add_text_box(slide, text, left, top, width, height, font_size=14, bold=False, color=None):
    textbox = slide.shapes.add_textbox(left, top, width, height)
    text_frame = textbox.text_frame
    text_frame.text = text
    text_frame.word_wrap = True

    for paragraph in text_frame.paragraphs:
        paragraph.font.size = Pt(font_size)
        paragraph.font.bold = bold
        if color:
            paragraph.font.color.rgb = color

    return textbox

# Extract all slides from HTML
slides_elements = tree.xpath('//div[@class="slide"]')

print(f"Found {len(slides_elements)} slides")

# Process each slide
for idx, slide_elem in enumerate(slides_elements):
    print(f"Processing slide {idx + 1}...")

    # Get section title (h2)
    h2_elements = slide_elem.xpath('.//h2[@class="section-title" or contains(@class, "text-5xl")]')
    title = clean_text(h2_elements[0].text_content()) if h2_elements else f"Slide {idx + 1}"

    # Create slide based on content
    if idx == 0:
        # First slide - title slide
        h1_elem = slide_elem.xpath('.//h1')
        main_title = clean_text(h1_elem[0].text_content()) if h1_elem else title

        p_elem = slide_elem.xpath('.//p[@class="text-2xl"]')
        subtitle = clean_text(p_elem[0].text_content()) if p_elem else ""

        slide = add_title_slide(main_title, subtitle)
    else:
        # Content slide
        slide = add_content_slide(title)

        # Check for images
        img_elements = slide_elem.xpath('.//img[contains(@src, "images/")]')

        # Layout content based on presence of images
        if img_elements:
            # Slide with image
            for img_elem in img_elements:
                img_src = img_elem.get('src')
                if img_src:
                    img_path = os.path.join('/Users/anasabounouar/Downloads/dbaichi/pfe-oracle', img_src)
                    alt_text = clean_text(img_elem.get('alt', ''))

                    # Add image centered
                    add_image_to_slide(slide, img_path, Inches(1.5), Inches(1.8), width=Inches(10))

        # Add text content
        # Get all paragraphs and list items
        p_elements = slide_elem.xpath('.//p[not(ancestor::div[contains(@class, "bg-")])]')
        li_elements = slide_elem.xpath('.//li')

        if p_elements or li_elements:
            # Add text box for content
            text_top = Inches(5.5) if img_elements else Inches(1.8)
            textbox = slide.shapes.add_textbox(Inches(0.5), text_top, Inches(12), Inches(1.5))
            text_frame = textbox.text_frame
            text_frame.word_wrap = True

            # Add paragraphs
            for p_idx, p_elem in enumerate(p_elements[:5]):  # Limit to 5 paragraphs
                text = clean_text(p_elem.text_content())
                if text and len(text) > 10:  # Skip very short text
                    if p_idx > 0:
                        text_frame.add_paragraph()
                    p = text_frame.paragraphs[-1]
                    p.text = text
                    p.font.size = Pt(12)

            # Add list items
            for li_elem in li_elements[:10]:  # Limit to 10 items
                text = clean_text(li_elem.text_content())
                if text:
                    text_frame.add_paragraph()
                    p = text_frame.paragraphs[-1]
                    p.text = "• " + text
                    p.font.size = Pt(11)
                    p.level = 0

# Save presentation
output_file = '/Users/anasabounouar/Downloads/dbaichi/pfe-oracle/presentation.pptx'
prs.save(output_file)
print(f"\nPresentation saved to: {output_file}")
print(f"Total slides created: {len(prs.slides)}")
