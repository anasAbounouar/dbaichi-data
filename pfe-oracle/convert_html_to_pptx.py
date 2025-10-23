#!/usr/bin/env python3
"""
HTML to PowerPoint Converter for GraalVM Infrastructure Presentation
Converts the presentation.html file into a professional PowerPoint presentation
preserving the visual design and content structure.
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup
from pathlib import Path
import re
import os

# Color scheme extracted from CSS
PRIMARY_COLOR = RGBColor(20, 83, 95)  # #14535F
ACCENT_COLOR = RGBColor(199, 70, 52)   # #C74634
TEAL_COLOR = RGBColor(15, 118, 110)    # #0f766e
BACKGROUND_COLOR = RGBColor(248, 250, 252)  # #f8fafc
WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(51, 65, 85)       # #334155
GRAY_LIGHT = RGBColor(100, 116, 139)   # #64748b

class HTMLToPowerPointConverter:
    """Converts HTML presentation to PowerPoint format"""

    def __init__(self, html_path, output_path, images_dir):
        self.html_path = Path(html_path)
        self.output_path = Path(output_path)
        self.images_dir = Path(images_dir)
        self.prs = Presentation()
        self.prs.slide_width = Inches(13.333)  # 16:9 aspect ratio
        self.prs.slide_height = Inches(7.5)

    def parse_html(self):
        """Parse the HTML file and extract slide content"""
        with open(self.html_path, 'r', encoding='utf-8') as f:
            html_content = f.read()

        self.soup = BeautifulSoup(html_content, 'html.parser')
        self.slides_data = []

        # Find all slide divs
        slides = self.soup.find_all('div', class_='slide')

        for slide in slides:
            slide_data = {
                'classes': slide.get('class', []),
                'content': slide,
                'is_divider': 'divider-slide' in slide.get('class', []),
                'is_title': 'title-slide' in slide.get('class', []),
            }
            self.slides_data.append(slide_data)

        return self.slides_data

    def clean_text(self, text):
        """Clean and normalize text content"""
        if not text:
            return ""
        # Remove extra whitespace
        text = re.sub(r'\s+', ' ', text)
        return text.strip()

    def add_title_slide(self, slide_data):
        """Create the title slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout

        # Set gradient background (approximation with solid color)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR

        # Add logos
        logo_ehtp_path = self.images_dir / "logo_ehtp.jpg"
        logo_graalvm_path = self.images_dir / "logo_graalvm.png"

        left_margin = Inches(0.5)
        logo_top = Inches(0.5)

        if logo_ehtp_path.exists():
            slide.shapes.add_picture(str(logo_ehtp_path), left_margin, logo_top, height=Inches(1.2))

        if logo_graalvm_path.exists():
            slide.shapes.add_picture(str(logo_graalvm_path), Inches(11.5), logo_top, height=Inches(1.2))

        # Add main title
        title_box = slide.shapes.add_textbox(Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.5))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True

        title_text = title_frame.paragraphs[0]
        title_text.text = "Infrastructure as Code Support\nin Graal CI"
        title_text.font.size = Pt(54)
        title_text.font.bold = True
        title_text.font.color.rgb = WHITE
        title_text.alignment = PP_ALIGN.CENTER

        # Add subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(4.2), Inches(9.3), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.text = "JSON-Driven IaC Framework with Pulumi & GitLab CI/CD"
        subtitle_para.font.size = Pt(24)
        subtitle_para.font.color.rgb = WHITE
        subtitle_para.alignment = PP_ALIGN.CENTER

        # Add author info
        author_box = slide.shapes.add_textbox(Inches(3), Inches(5.5), Inches(7.3), Inches(1.5))
        author_frame = author_box.text_frame
        author_frame.word_wrap = True

        lines = [
            "Yassine DBAICHI",
            "École Hassania des Travaux Publics",
            "Oracle Labs - GraalVM RISQ Team",
            "Academic Year 2024-2025"
        ]

        for i, line in enumerate(lines):
            if i > 0:
                author_frame.add_paragraph()
            para = author_frame.paragraphs[i]
            para.text = line
            para.font.size = Pt(18) if i == 0 else Pt(16)
            para.font.bold = (i == 0)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER
            para.space_after = Pt(6)

    def add_divider_slide(self, title_text):
        """Create a section divider slide"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR

        # Add title
        title_box = slide.shapes.add_textbox(Inches(1), Inches(2.5), Inches(11.3), Inches(2))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        title_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

        title_para = title_frame.paragraphs[0]
        title_para.text = title_text
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = WHITE
        title_para.alignment = PP_ALIGN.CENTER

    def add_content_slide(self, slide_data):
        """Create a content slide with bullet points"""
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BACKGROUND_COLOR

        content = slide_data['content']

        # Extract title
        title_elem = content.find('h2', class_='slide-title')
        if title_elem:
            title_text = self.clean_text(title_elem.get_text())

            # Add title with accent underline
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12.3), Inches(0.8))
            title_frame = title_box.text_frame
            title_para = title_frame.paragraphs[0]
            title_para.text = title_text
            title_para.font.size = Pt(36)
            title_para.font.bold = True
            title_para.font.color.rgb = PRIMARY_COLOR

            # Add accent line
            line = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(0.5), Inches(1.15),
                Inches(12.3), Inches(0.05)
            )
            line.fill.solid()
            line.fill.fore_color.rgb = ACCENT_COLOR
            line.line.fill.background()

        # Extract content
        slide_content = content.find('div', class_='slide-content')

        if slide_content:
            # Check for special layouts
            if slide_content.find('div', class_='card-grid'):
                self._add_card_grid_content(slide, slide_content)
            elif slide_content.find('div', class_='two-column'):
                self._add_two_column_content(slide, slide_content)
            elif slide_content.find('ul', class_='bullet-list'):
                self._add_bullet_list_content(slide, slide_content)
            elif slide_content.find('img'):
                self._add_image_content(slide, slide_content)
            else:
                # Generic content
                self._add_generic_content(slide, slide_content)

    def _add_bullet_list_content(self, slide, content_div):
        """Add bullet list content to slide"""
        bullet_list = content_div.find('ul', class_='bullet-list')
        if not bullet_list:
            return

        # Create text box for bullets
        text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.4), Inches(11.9), Inches(5.6))
        text_frame = text_box.text_frame
        text_frame.word_wrap = True

        items = bullet_list.find_all('li', recursive=False)

        for idx, item in enumerate(items):
            if idx > 0:
                text_frame.add_paragraph()

            para = text_frame.paragraphs[idx]

            # Extract text, handling <strong> tags
            item_text = ""
            for child in item.children:
                if child.name == 'strong':
                    item_text += child.get_text()
                elif child.name is None:  # Text node
                    item_text += child
                elif child.name == 'ul':
                    # Skip nested lists for now
                    pass
                else:
                    item_text += child.get_text()

            item_text = self.clean_text(item_text)
            para.text = item_text
            para.level = 0
            para.font.size = Pt(18)
            para.font.color.rgb = TEXT_DARK
            para.space_before = Pt(8)
            para.space_after = Pt(8)

        # Add highlight box if present
        highlight = content_div.find('div', class_='highlight-box')
        if highlight:
            highlight_text = self.clean_text(highlight.get_text())

            # Create highlight box
            box = slide.shapes.add_shape(
                1,  # Rectangle
                Inches(1.5), Inches(6.3),
                Inches(10.3), Inches(0.8)
            )
            box.fill.solid()
            box.fill.fore_color.rgb = PRIMARY_COLOR
            box.line.fill.background()

            # Add highlight text
            text_frame = box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.2)
            text_frame.margin_right = Inches(0.2)
            text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE

            para = text_frame.paragraphs[0]
            para.text = highlight_text
            para.font.size = Pt(16)
            para.font.color.rgb = WHITE
            para.alignment = PP_ALIGN.CENTER

    def _add_card_grid_content(self, slide, content_div):
        """Add card grid layout to slide"""
        card_grid = content_div.find('div', class_='card-grid')
        cards = card_grid.find_all('div', class_='card')

        # Calculate card positions (2x3 grid)
        card_width = Inches(5.5)
        card_height = Inches(1.8)
        h_spacing = Inches(0.5)
        v_spacing = Inches(0.3)
        start_left = Inches(0.7)
        start_top = Inches(1.5)

        for idx, card in enumerate(cards[:6]):  # Max 6 cards
            row = idx // 2
            col = idx % 2

            left = start_left + col * (card_width + h_spacing)
            top = start_top + row * (card_height + v_spacing)

            # Create card shape
            card_shape = slide.shapes.add_shape(
                1,  # Rectangle
                left, top, card_width, card_height
            )
            card_shape.fill.solid()
            card_shape.fill.fore_color.rgb = WHITE
            card_shape.line.color.rgb = PRIMARY_COLOR
            card_shape.line.width = Pt(2)

            # Add card content
            text_frame = card_shape.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_right = Inches(0.15)
            text_frame.margin_top = Inches(0.1)
            text_frame.margin_bottom = Inches(0.1)

            # Card title
            card_title = card.find('div', class_='card-title')
            if card_title:
                title_text = self.clean_text(card_title.get_text())
                para = text_frame.paragraphs[0]
                para.text = title_text
                para.font.size = Pt(16)
                para.font.bold = True
                para.font.color.rgb = PRIMARY_COLOR
                para.space_after = Pt(6)

            # Card content
            card_content = card.find('div', class_='card-content')
            if card_content:
                content_text = self.clean_text(card_content.get_text())
                text_frame.add_paragraph()
                para = text_frame.paragraphs[1]
                para.text = content_text
                para.font.size = Pt(13)
                para.font.color.rgb = TEXT_DARK
                para.line_spacing = 1.2

    def _add_two_column_content(self, slide, content_div):
        """Add two-column layout to slide"""
        two_col = content_div.find('div', class_='two-column')
        columns = two_col.find_all(['div'], recursive=False)

        if len(columns) >= 2:
            # Left column
            left_col = columns[0]
            self._add_column_content(slide, left_col, Inches(0.7), Inches(1.5), Inches(5.6))

            # Right column
            right_col = columns[1]
            # Check if it contains an image
            if right_col.find('img'):
                self._add_image_to_column(slide, right_col, Inches(7.2), Inches(1.5), Inches(5.6))
            else:
                self._add_column_content(slide, right_col, Inches(7.2), Inches(1.5), Inches(5.6))

    def _add_column_content(self, slide, column, left, top, width):
        """Add content to a column"""
        # Check for bullet list
        bullet_list = column.find('ul', class_='bullet-list')
        if bullet_list:
            text_box = slide.shapes.add_textbox(left, top, width, Inches(5.5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            items = bullet_list.find_all('li', recursive=False)
            for idx, item in enumerate(items):
                if idx > 0:
                    text_frame.add_paragraph()

                para = text_frame.paragraphs[idx]
                item_text = self.clean_text(item.get_text())
                para.text = item_text
                para.level = 0
                para.font.size = Pt(16)
                para.font.color.rgb = TEXT_DARK
                para.space_after = Pt(6)

        # Check for cards
        cards = column.find_all('div', class_='card')
        if cards:
            current_top = top
            for card in cards:
                card_height = Inches(1.5)

                # Create card shape
                card_shape = slide.shapes.add_shape(
                    1,  # Rectangle
                    left, current_top, width, card_height
                )
                card_shape.fill.solid()
                card_shape.fill.fore_color.rgb = WHITE
                card_shape.line.color.rgb = PRIMARY_COLOR
                card_shape.line.width = Pt(2)

                text_frame = card_shape.text_frame
                text_frame.word_wrap = True
                text_frame.margin_left = Inches(0.15)
                text_frame.margin_right = Inches(0.15)
                text_frame.margin_top = Inches(0.1)

                # Card title
                card_title = card.find('div', class_='card-title')
                if card_title:
                    title_text = self.clean_text(card_title.get_text())
                    para = text_frame.paragraphs[0]
                    para.text = title_text
                    para.font.size = Pt(14)
                    para.font.bold = True
                    para.font.color.rgb = PRIMARY_COLOR

                # Card content
                card_content = card.find('div', class_='card-content')
                if card_content:
                    content_text = self.clean_text(card_content.get_text())
                    text_frame.add_paragraph()
                    para = text_frame.paragraphs[1]
                    para.text = content_text
                    para.font.size = Pt(12)
                    para.font.color.rgb = TEXT_DARK

                current_top += card_height + Inches(0.15)

        # Check for highlight box
        highlight = column.find('div', class_='highlight-box')
        if highlight:
            highlight_text = self.clean_text(highlight.get_text())

            # Check if there's a title (h3)
            h3_title = highlight.find('h3')
            title_text = ""
            if h3_title:
                title_text = self.clean_text(h3_title.get_text())
                # Remove title from highlight text
                highlight_text = highlight_text.replace(title_text, "").strip()

            box_height = Inches(2.5)
            box = slide.shapes.add_shape(
                1,  # Rectangle
                left, Inches(4.3), width, box_height
            )
            box.fill.solid()
            box.fill.fore_color.rgb = PRIMARY_COLOR
            box.line.fill.background()

            text_frame = box.text_frame
            text_frame.word_wrap = True
            text_frame.margin_left = Inches(0.15)
            text_frame.margin_right = Inches(0.15)
            text_frame.margin_top = Inches(0.15)

            if title_text:
                para = text_frame.paragraphs[0]
                para.text = title_text
                para.font.size = Pt(18)
                para.font.bold = True
                para.font.color.rgb = WHITE
                para.space_after = Pt(8)

                text_frame.add_paragraph()
                para = text_frame.paragraphs[1]
                para.text = highlight_text
                para.font.size = Pt(14)
                para.font.color.rgb = WHITE
            else:
                para = text_frame.paragraphs[0]
                para.text = highlight_text
                para.font.size = Pt(14)
                para.font.color.rgb = WHITE

    def _add_image_to_column(self, slide, column, left, top, width):
        """Add image to a column"""
        img_container = column.find('div', class_='image-container')
        if img_container:
            img_elem = img_container.find('img')
            if img_elem:
                img_src = img_elem.get('src', '')
                # Handle relative path
                if img_src.startswith('images/'):
                    img_path = self.images_dir / img_src.replace('images/', '')
                else:
                    img_path = self.images_dir / img_src

                if img_path.exists():
                    try:
                        slide.shapes.add_picture(
                            str(img_path),
                            left, top,
                            width=width,
                            height=Inches(5)
                        )
                    except Exception as e:
                        print(f"Warning: Could not add image {img_path}: {e}")

    def _add_image_content(self, slide, content_div):
        """Add full-width image content"""
        img_elem = content_div.find('img')
        if img_elem:
            img_src = img_elem.get('src', '')
            # Handle relative path
            if img_src.startswith('images/'):
                img_path = self.images_dir / img_src.replace('images/', '')
            else:
                img_path = self.images_dir / img_src

            if img_path.exists():
                try:
                    # Center the image
                    slide.shapes.add_picture(
                        str(img_path),
                        Inches(1.5), Inches(1.5),
                        width=Inches(10.3)
                    )
                except Exception as e:
                    print(f"Warning: Could not add image {img_path}: {e}")

    def _add_generic_content(self, slide, content_div):
        """Add generic text content"""
        text_content = self.clean_text(content_div.get_text())

        if text_content:
            text_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.5), Inches(11.9), Inches(5.5))
            text_frame = text_box.text_frame
            text_frame.word_wrap = True

            para = text_frame.paragraphs[0]
            para.text = text_content
            para.font.size = Pt(18)
            para.font.color.rgb = TEXT_DARK

    def convert(self):
        """Main conversion process"""
        print("Parsing HTML file...")
        self.parse_html()

        print(f"Found {len(self.slides_data)} slides")

        for idx, slide_data in enumerate(self.slides_data):
            print(f"Processing slide {idx + 1}/{len(self.slides_data)}...")

            if slide_data['is_title']:
                self.add_title_slide(slide_data)
            elif slide_data['is_divider']:
                # Extract divider title
                divider_title = slide_data['content'].find('h2', class_='divider-title')
                if divider_title:
                    title_text = self.clean_text(divider_title.get_text())
                    self.add_divider_slide(title_text)
            else:
                # Content slide
                self.add_content_slide(slide_data)

        print(f"Saving presentation to {self.output_path}...")
        self.prs.save(str(self.output_path))
        print(f"Conversion complete! Created {len(self.prs.slides)} slides.")


def main():
    """Main entry point"""
    # Set paths
    base_dir = Path("/Users/anasabounouar/Downloads/dbaichi/pfe-oracle")
    html_path = base_dir / "presentation.html"
    output_path = base_dir / "presentation_converted.pptx"
    images_dir = base_dir / "images"

    # Verify files exist
    if not html_path.exists():
        print(f"Error: HTML file not found at {html_path}")
        return

    if not images_dir.exists():
        print(f"Warning: Images directory not found at {images_dir}")

    # Create converter and run
    converter = HTMLToPowerPointConverter(html_path, output_path, images_dir)
    converter.convert()

    print(f"\nPowerPoint presentation created successfully!")
    print(f"Output file: {output_path}")


if __name__ == "__main__":
    main()
